from pptx import Presentation
import os

def process_ppt(ppt_path):
    prs = Presentation(ppt_path)
    full_text = []
    
    for i, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        full_text.append(f"--- Slide {i+1} ---\n" + "\n".join(slide_text))
        
    return "\n\n".join(full_text)
